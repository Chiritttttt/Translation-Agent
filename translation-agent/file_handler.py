import os
import base64
import sys

# ─── OCR 进度回调 ────────────────────────────────────
_ocr_progress_callback = None


def set_ocr_progress_callback(callback):
    """设置 OCR 进度回调函数，供 GUI 调用。
    callback(current_page, total_pages, message)
    """
    global _ocr_progress_callback
    _ocr_progress_callback = callback


def _report_ocr_progress(current, total, msg):
    if _ocr_progress_callback:
        try:
            _ocr_progress_callback(current, total, msg)
        except Exception:
            pass


# ─── 读取 ────────────────────────────────────────────────

def read_pdf(file_path):
    """读取 PDF 文件，自动判断是否为扫描件并启用 AI OCR。

    流程：
    1. 用 PyMuPDF 提取文本
    2. 质量检测：平均每页 < 20 字符 → 判定为扫描件
    3. 扫描件 → pdf2image 转图片 → OpenAI 兼容 Vision API 做 OCR
    """
    try:
        import fitz
        doc = fitz.open(file_path)
        page_count = len(doc)

        # 第一步：尝试文本提取
        text_parts = []
        for page in doc:
            page_text = page.get_text().strip()
            text_parts.append(page_text)
        full_text = "\n\n".join(text_parts).strip()

        # 第二步：质量检测
        avg_chars = len(full_text) / max(page_count, 1)
        if avg_chars >= 20:
            return full_text

        # 第三步：文本太少，很可能是扫描件，启用 OCR
        _report_ocr_progress(0, page_count, "检测到扫描件，启用 AI 视觉 OCR...")
        ocr_text = _ocr_with_vision(file_path, page_count)
        if ocr_text and len(ocr_text.strip()) > 10:
            return ocr_text.strip()

        # OCR 也失败，返回原始文本（可能为空）
        return full_text

    except Exception as e:
        raise Exception(f"PDF读取失败: {e}")


def _ocr_with_vision(file_path, page_count):
    """使用 OpenAI 兼容的 Vision API 对 PDF 逐页做 OCR。
    读取 .env 中的 OPENAI_API_KEY / OPENAI_BASE_URL / OPENAI_MODEL 配置。
    """
    try:
        from dotenv import load_dotenv
        load_dotenv()

        import openai
        from pdf2image import convert_from_path
        from io import BytesIO

        api_key = os.getenv("OPENAI_API_KEY", "")
        base_url = os.getenv("OPENAI_BASE_URL", "")
        model = os.getenv("OPENAI_MODEL", "")

        if not api_key:
            print("[OCR] 未配置 OPENAI_API_KEY，无法进行 OCR 识别")
            return ""

        # 选择模型：优先使用用户配置的模型，带 vision 能力的模型
        vision_model = model if model else "gpt-4o"

        client = openai.OpenAI(api_key=api_key, base_url=base_url if base_url else None)

        # 将 PDF 每页转为图片
        _report_ocr_progress(0, page_count, "正在转换 PDF 页面为图片...")
        images = convert_from_path(file_path, dpi=200)

        if not images:
            return ""

        actual_pages = len(images)
        all_text = []

        for i, img in enumerate(images):
            _report_ocr_progress(i + 1, actual_pages,
                                 f"OCR 识别中: 第 {i+1}/{actual_pages} 页...")
            try:
                # 将 PIL Image 转 base64
                buf = BytesIO()
                # 统一转为 RGB 避免 RGBA 问题
                if img.mode == "RGBA":
                    img = img.convert("RGB")
                img.save(buf, format="PNG")
                img_b64 = base64.b64encode(buf.getvalue()).decode("utf-8")

                # 调用 Vision API
                response = client.chat.completions.create(
                    model=vision_model,
                    messages=[
                        {
                            "role": "system",
                            "content": (
                                "你是一个专业的 OCR 助手。请从文档图片中完整、准确地提取所有文字。"
                                "保留段落、标题、列表等结构。只输出提取的文字，不要有任何解释。"
                            ),
                        },
                        {
                            "role": "user",
                            "content": [
                                {
                                    "type": "text",
                                    "text": f"请提取第 {i+1} 页的所有文字，保留原始格式和结构。",
                                },
                                {
                                    "type": "image_url",
                                    "image_url": {
                                        "url": f"data:image/png;base64,{img_b64}"
                                    },
                                },
                            ],
                        },
                    ],
                    temperature=0.1,
                    max_tokens=4096,
                )

                page_text = response.choices[0].message.content or ""
                if page_text.strip():
                    all_text.append(page_text.strip())

            except Exception as e:
                print(f"[OCR] 第 {i+1} 页识别失败: {e}")
                continue

        _report_ocr_progress(actual_pages, actual_pages, "OCR 识别完成！")
        return "\n\n".join(all_text)

    except ImportError as e:
        print(f"[OCR] 缺少依赖: {e}")
        print("[OCR] 请安装: pip install pdf2image Pillow openai python-dotenv")
        return ""
    except Exception as e:
        print(f"[OCR] Vision API 调用失败: {e}")
        return ""


def read_docx(file_path):
    try:
        from docx import Document
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text.strip()
    except Exception as e:
        raise Exception(f"DOCX读取失败: {e}")


def read_pptx(file_path):
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n\n"
        return text.strip()
    except Exception as e:
        raise Exception(f"PPTX读取失败: {e}")


def read_excel(file_path):
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text = ""
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                row_text = "  ".join([str(c) for c in row if c is not None])
                if row_text.strip():
                    text += row_text + "\n"
        return text.strip()
    except Exception as e:
        raise Exception(f"Excel读取失败: {e}")


def read_file(file_path):
    """读取文件，自动识别格式。

    返回: (file_type, content, extra)
      - file_type: 'pdf' | 'docx' | 'pptx' | 'xlsx' | 'txt'
      - content: 提取的文本
      - extra: Excel 时为行数据列表 [[cell, ...], ...]，其他格式为 None
    """
    ext = file_path.lower().split(".")[-1]
    if ext == "pdf":
        return "pdf", read_pdf(file_path), None
    elif ext in ["docx", "doc"]:
        return "docx", read_docx(file_path), None
    elif ext in ["pptx", "ppt"]:
        return "pptx", read_pptx(file_path), None
    elif ext in ["xlsx", "xls"]:
        return "xlsx", read_excel(file_path), None
    else:
        raise Exception("不支持此格式")


# ─── 工具：译文行对齐 ────────────────────────────────────

def match_translation(original_lines, translated_text):
    trans_lines = [l for l in translated_text.split("\n") if l.strip()]
    orig_count = len(original_lines)
    trans_count = len(trans_lines)
    if orig_count == 0:
        return []
    if trans_count == 0:
        return [""] * orig_count
    result = []
    for i in range(orig_count):
        idx = min(int(i * trans_count / orig_count), trans_count - 1)
        result.append(trans_lines[idx])
    return result


# ─── 工具：复制 DOCX 格式 ────────────────────────────────

def _copy_run_format(src_run, dst_run):
    try:
        dst_run.bold = src_run.bold
        dst_run.italic = src_run.italic
        dst_run.underline = src_run.underline
        if src_run.font.size:
            dst_run.font.size = src_run.font.size
        if src_run.font.color and src_run.font.color.type:
            try:
                dst_run.font.color.rgb = src_run.font.color.rgb
            except Exception:
                pass
        if src_run.font.name:
            dst_run.font.name = src_run.font.name
    except Exception:
        pass


def _copy_para_format(src_para, dst_para):
    try:
        dst_para.alignment = src_para.alignment
        try:
            dst_para.style = src_para.style
        except Exception:
            pass
        pf = src_para.paragraph_format
        dpf = dst_para.paragraph_format
        if pf.space_before:
            dpf.space_before = pf.space_before
        if pf.space_after:
            dpf.space_after = pf.space_after
        if pf.line_spacing:
            dpf.line_spacing = pf.line_spacing
        if pf.left_indent:
            dpf.left_indent = pf.left_indent
        if pf.first_line_indent:
            dpf.first_line_indent = pf.first_line_indent
    except Exception:
        pass


# ─── 导出 DOCX ───────────────────────────────────────────

def export_docx_translation(translated, output_path, original_path=None):
    from docx import Document
    if original_path and os.path.exists(original_path):
        src_doc = Document(original_path)
        dst_doc = Document(original_path)
        orig_paras = [p for p in src_doc.paragraphs if p.text.strip()]
        matched = match_translation([p.text for p in orig_paras], translated)
        trans_idx = 0
        for para in dst_doc.paragraphs:
            if para.text.strip() and trans_idx < len(matched):
                new_text = matched[trans_idx]
                trans_idx += 1
                if para.runs:
                    para.runs[0].text = new_text
                    for run in para.runs[1:]:
                        run.text = ""
                else:
                    para.add_run(new_text)
        dst_doc.save(output_path)
    else:
        doc = Document()
        for line in translated.split("\n"):
            if line.strip():
                doc.add_paragraph(line)
        doc.save(output_path)


def export_docx_paragraph(original, translated, output_path, original_path=None):
    """段落对照：一段原文一段译文，完全复制原文格式，无装饰"""
    from docx import Document
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    import copy

    if original_path and os.path.exists(original_path):
        src_doc = Document(original_path)
        orig_src_paras = [p for p in src_doc.paragraphs if p.text.strip()]
    else:
        orig_src_paras = None

    # 从原文文件直接复制作为基础
    if original_path and os.path.exists(original_path):
        dst_doc = Document(original_path)
        # 清空所有段落内容
        for para in dst_doc.paragraphs:
            for run in para.runs:
                run.text = ""
    else:
        dst_doc = Document()

    raw_orig = [p for p in original.split("\n") if p.strip()]
    matched_trans = match_translation(raw_orig, translated)

    # 清除文档内容，重新写入
    from docx.oxml.ns import qn
    body = dst_doc.element.body
    # 保留最后一个sectPr（页面设置），清除其余内容
    sect_pr = body.find(qn('w:sectPr'))
    for child in list(body):
        if child != sect_pr:
            body.remove(child)

    for i, orig_text in enumerate(raw_orig):
        # ── 原文段落：完整复制格式 ──
        orig_p = dst_doc.add_paragraph()
        if orig_src_paras and i < len(orig_src_paras):
            _copy_para_format(orig_src_paras[i], orig_p)
            for src_run in orig_src_paras[i].runs:
                dst_run = orig_p.add_run(src_run.text)
                _copy_run_format(src_run, dst_run)
        else:
            orig_p.add_run(orig_text)

        # ── 译文段落：完全复制原文格式，只换文字 ──
        trans_text = matched_trans[i] if i < len(matched_trans) else ""
        trans_p = dst_doc.add_paragraph()
        if orig_src_paras and i < len(orig_src_paras):
            _copy_para_format(orig_src_paras[i], trans_p)
            if orig_src_paras[i].runs:
                # 只用第一个run的格式，文字换成译文
                dst_run = trans_p.add_run(trans_text)
                _copy_run_format(orig_src_paras[i].runs[0], dst_run)
            else:
                trans_p.add_run(trans_text)
        else:
            trans_p.add_run(trans_text)

    dst_doc.save(output_path)


def export_docx_bilingual(original, translated, output_path, original_path=None):
    """行对照：双列表格，去掉标题"""
    from docx import Document
    from docx.shared import RGBColor
    from docx.enum.table import WD_TABLE_ALIGNMENT

    if original_path and os.path.exists(original_path):
        src_doc = Document(original_path)
        orig_src_paras = [p for p in src_doc.paragraphs if p.text.strip()]
    else:
        orig_src_paras = None

    doc = Document()
    # ← 去掉了 add_heading
    raw_orig = [p for p in original.split("\n") if p.strip()]
    matched_trans = match_translation(raw_orig, translated)

    table = doc.add_table(rows=1, cols=2)
    table.style = "Table Grid"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    hdr = table.rows[0].cells
    for idx, txt in enumerate(["原文", "译文"]):
        run = hdr[idx].paragraphs[0].add_run(txt)
        run.bold = True
        run.font.color.rgb = RGBColor(0x4F, 0x8E, 0xF7)

    for i, orig_text in enumerate(raw_orig):
        row = table.add_row().cells
        orig_cell_para = row[0].paragraphs[0]
        if orig_src_paras and i < len(orig_src_paras):
            _copy_para_format(orig_src_paras[i], orig_cell_para)
            for src_run in orig_src_paras[i].runs:
                dst_run = orig_cell_para.add_run(src_run.text)
                _copy_run_format(src_run, dst_run)
        else:
            orig_cell_para.add_run(orig_text)

        trans_cell_para = row[1].paragraphs[0]
        trans_text = matched_trans[i] if i < len(matched_trans) else ""
        if orig_src_paras and i < len(orig_src_paras):
            _copy_para_format(orig_src_paras[i], trans_cell_para)
            if orig_src_paras[i].runs:
                dst_run = trans_cell_para.add_run(trans_text)
                _copy_run_format(orig_src_paras[i].runs[0], dst_run)
            else:
                trans_cell_para.add_run(trans_text)
        else:
            trans_cell_para.add_run(trans_text)

    doc.save(output_path)


# ─── 导出 PPTX ───────────────────────────────────────────

def export_pptx_translation(original_path, translated, output_path):
    """全译文：保留格式只换文字"""
    from pptx import Presentation
    prs = Presentation(original_path)
    trans_lines = [l for l in translated.split("\n") if l.strip()
                   and not (l.startswith("[幻灯片") and l.endswith("]"))]
    line_idx = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                if para.text.strip():
                    new_text = trans_lines[line_idx] if line_idx < len(trans_lines) else ""
                    line_idx += 1
                    if para.runs:
                        para.runs[0].text = new_text
                        for run in para.runs[1:]:
                            run.text = ""
    prs.save(output_path)


def export_pptx_bilingual(original_path, translated, output_path):
    """备注栏对照：正文保留原文，译文写入备注"""
    from pptx import Presentation
    prs = Presentation(original_path)
    trans_lines = [l for l in translated.split("\n") if l.strip()
                   and not (l.startswith("[幻灯片") and l.endswith("]"))]
    line_idx = 0
    for slide in prs.slides:
        slide_trans = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        t = trans_lines[line_idx] if line_idx < len(trans_lines) else ""
                        slide_trans.append(t)
                        line_idx += 1
        if slide_trans:
            notes = slide.notes_slide
            notes.notes_text_frame.text = "【译文】\n" + "\n".join(slide_trans)
    prs.save(output_path)


# ─── 导出 Excel ──────────────────────────────────────────

def export_excel_bilingual(original_rows, translated, output_path, original_path=None):
    import openpyxl
    import copy
    from openpyxl.styles import PatternFill, Font

    if original_path and os.path.exists(original_path):
        wb = openpyxl.load_workbook(original_path)
    else:
        wb = openpyxl.Workbook()

    ws = wb.active
    trans_rows = [r for r in translated.split("\n") if r.strip()]
    max_col = ws.max_column

    header = ws.cell(row=1, column=max_col + 1, value="译文")
    header.fill = PatternFill("solid", fgColor="4F8EF7")
    header.font = Font(color="FFFFFF", bold=True)

    for i in range(2, ws.max_row + 1):
        trans_text = trans_rows[i - 2] if (i - 2) < len(trans_rows) else ""
        src_cell = ws.cell(row=i, column=max_col)
        new_cell = ws.cell(row=i, column=max_col + 1, value=trans_text)
        try:
            new_cell.font = copy.copy(src_cell.font)
            new_cell.alignment = copy.copy(src_cell.alignment)
            new_cell.border = copy.copy(src_cell.border)
        except Exception:
            pass

    wb.save(output_path)


def export_excel_translation(translated, output_path, original_path=None):
    import openpyxl
    import copy

    if original_path and os.path.exists(original_path):
        wb = openpyxl.load_workbook(original_path)
        ws = wb.active
        trans_rows = [r for r in translated.split("\n") if r.strip()]
        row_idx = 0
        for row in ws.iter_rows():
            if any(cell.value is not None for cell in row):
                row_text = trans_rows[row_idx] if row_idx < len(trans_rows) else ""
                cells_text = row_text.split("  ")
                for c_idx, cell in enumerate(row):
                    if cell.value is not None:
                        cell.value = cells_text[c_idx] if c_idx < len(cells_text) else ""
                row_idx += 1
    else:
        wb = openpyxl.Workbook()
        ws = wb.active
        for i, line in enumerate(translated.split("\n"), start=1):
            if line.strip():
                for j, val in enumerate(line.split("  "), start=1):
                    ws.cell(row=i, column=j, value=val)

    wb.save(output_path)


# ─── 导出 PDF ────────────────────────────────────────────

def register_fonts():
    try:
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        # 按优先级搜索 CJK 字体路径
        font_paths = [
            # Windows
            "C:/Windows/Fonts/msyh.ttc",
            "C:/Windows/Fonts/simsun.ttc",
            "C:/Windows/Fonts/simhei.ttf",
            # Linux
            "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
            "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
            "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",
            "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
            "/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf",
            "/usr/share/fonts/noto-cjk/NotoSansCJK-Regular.ttc",
            # macOS
            "/System/Library/Fonts/PingFang.ttc",
            "/System/Library/Fonts/STHeiti Medium.ttc",
        ]
        for path in font_paths:
            if os.path.exists(path):
                pdfmetrics.registerFont(TTFont("CJK", path))
                return "CJK"
    except Exception:
        pass
    return "Helvetica"


def export_pdf_translation(translated, output_path):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    font = register_fonts()
    doc = SimpleDocTemplate(output_path, pagesize=A4,
                            leftMargin=2*cm, rightMargin=2*cm,
                            topMargin=2*cm, bottomMargin=2*cm)
    style = ParagraphStyle("b", fontName=font, fontSize=11,
                           leading=18, wordWrap="CJK")
    story = []
    for para in translated.split("\n"):
        if para.strip():
            story.append(Paragraph(para, style))
            story.append(Spacer(1, 6))
    doc.build(story)


def export_pdf_paragraph(original, translated, output_path):
    """段落对照PDF：一段原文一段译文"""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
    from reportlab.lib import colors
    font = register_fonts()
    doc = SimpleDocTemplate(output_path, pagesize=A4,
                            leftMargin=2*cm, rightMargin=2*cm,
                            topMargin=2*cm, bottomMargin=2*cm)
    orig_style = ParagraphStyle(
        "orig", fontName=font, fontSize=10, leading=16,
        wordWrap="CJK", textColor=colors.HexColor("#555555"),
        backColor=colors.HexColor("#F5F5F5"),
        borderPadding=(6, 8, 6, 8),
    )
    trans_style = ParagraphStyle(
        "trans", fontName=font, fontSize=10, leading=16,
        wordWrap="CJK", textColor=colors.HexColor("#1A1A1A"),
        leftIndent=10,
        borderLeftWidth=3,
        borderLeftColor=colors.HexColor("#4F8EF7"),
        borderLeftPadding=8,
    )
    orig_paras = [p for p in original.split("\n") if p.strip()]
    matched_trans = match_translation(orig_paras, translated)
    story = []
    for i, orig_text in enumerate(orig_paras):
        story.append(Paragraph(orig_text, orig_style))
        story.append(Spacer(1, 4))
        trans_text = matched_trans[i] if i < len(matched_trans) else ""
        story.append(Paragraph(trans_text, trans_style))
        story.append(Spacer(1, 4))
        story.append(HRFlowable(width="100%", thickness=0.5,
                                color=colors.HexColor("#DDDDDD")))
        story.append(Spacer(1, 8))
    doc.build(story)


def export_pdf_bilingual(original, translated, output_path):
    """双列PDF"""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph
    from reportlab.lib import colors
    font = register_fonts()
    doc = SimpleDocTemplate(output_path, pagesize=A4,
                            leftMargin=1.5*cm, rightMargin=1.5*cm,
                            topMargin=2*cm, bottomMargin=2*cm)
    style = ParagraphStyle("b", fontName=font, fontSize=9,
                           leading=14, wordWrap="CJK")
    hstyle = ParagraphStyle("h", fontName=font, fontSize=10, leading=14,
                            textColor=colors.HexColor("#4F8EF7"), wordWrap="CJK")
    orig_paras = [p for p in original.split("\n") if p.strip()]
    matched_trans = match_translation(orig_paras, translated)
    data = [[Paragraph("原文", hstyle), Paragraph("译文", hstyle)]]
    for i, orig_text in enumerate(orig_paras):
        t = matched_trans[i] if i < len(matched_trans) else ""
        data.append([Paragraph(orig_text, style), Paragraph(t, style)])
    col_w = (A4[0] - 3*cm) / 2
    tbl = Table(data, colWidths=[col_w, col_w], repeatRows=1)
    tbl.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#EEF4FF")),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.lightgrey),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1),
         [colors.white, colors.HexColor("#F9F9F9")]),
    ]))
    doc.build([tbl])